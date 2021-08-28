from pptx import Presentation
prs = Presentation()
prs.save('test.pptx')
#新建一页幻灯片
slide= prs.slides.add_slide(prs.slide_layouts[1])

#编辑幻灯中的元素
body_shape= slide.shapes.placeholders
# body_shape为本页ppt中所有shapes
body_shape[0].text= 'hello'
body_shape[1].text= 'world'

#添加表格
rows, cols, left, top, width, height= 2, 2, Inches(3.5), Inches(4.5), Inches(6), Inches(0.8)
table= slide.shapes.add_table(rows, cols, left, top, width, height).table
# 添加表格，并取表格类
table.columns[0].width= Inches(2.0)
# 第一纵列宽度
table.columns[1].width= Inches(4.0)
# 第二纵列宽度
table.cell(0, 0).text= 'table00'
# 指定位置写入文本
table.cell(0, 1).text= 'table01'
table.cell(1, 0).text= 'table10'
table.cell(1, 1).text= 'table11'